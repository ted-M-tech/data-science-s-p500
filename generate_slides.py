#!/usr/bin/env python3
"""Generate one slide deck per part (4 total), 8-10 slides each. Plain style."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

ROOT = Path(__file__).parent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)


def add_text(slide, left, top, w, h, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return tf


def title_slide(prs, title, subtitle):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(sl, 1, 2.2, 8, 1, title, size=32, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, 1, 3.5, 8, 0.8, subtitle, size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(sl, 1, 4.5, 8, 0.5, "Dataset: S&P 500 Stocks (Kaggle)", size=14, color=GRAY, align=PP_ALIGN.CENTER)


def content_slide(prs, title, bullets, part_dir=None, figure=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(sl, 0.5, 0.3, 9, 0.6, title, size=24, bold=True)

    if figure and part_dir:
        fig_path = ROOT / part_dir / 'figures' / figure
        if fig_path.exists():
            tf = add_text(sl, 0.5, 1.1, 4.5, 5.5, '', size=14)
            for i, b in enumerate(bullets):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = f"- {b}"
                p.font.size = Pt(14)
                p.font.color.rgb = DARK
                p.font.name = 'Calibri'
                p.space_after = Pt(6)
            sl.shapes.add_picture(str(fig_path), Inches(5.2), Inches(1.1), Inches(4.4))
            return

    tf = add_text(sl, 0.5, 1.1, 9, 5.5, '', size=16)
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"- {b}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.font.name = 'Calibri'
        p.space_after = Pt(8)


def table_slide(prs, title, headers, rows):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(sl, 0.5, 0.3, 9, 0.6, title, size=24, bold=True)
    cols = len(headers)
    tbl = sl.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.35 + 0.35 * len(rows))).table
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.name = 'Calibri'
    for r, row in enumerate(rows):
        for c_i, v in enumerate(row):
            c = tbl.cell(r + 1, c_i)
            c.text = str(v)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = 'Calibri'


def figure_slide(prs, title, part_dir, figure, caption=''):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(sl, 0.5, 0.3, 9, 0.6, title, size=24, bold=True)
    fig_path = ROOT / part_dir / 'figures' / figure
    if fig_path.exists():
        sl.shapes.add_picture(str(fig_path), Inches(1.5), Inches(1.2), Inches(7))
    if caption:
        add_text(sl, 1, 6.2, 8, 0.5, caption, size=12, color=GRAY, align=PP_ALIGN.CENTER)


def build_part1():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    d = 'part1_regression'

    # 1. Title
    title_slide(prs, "Part 1: Regression", "Predicting S&P 500 Stock Closing Prices")

    # 2. Dataset
    content_slide(prs, "Dataset", [
        "Source: S&P 500 Stocks (Kaggle) -- daily OHLCV data",
        "Top 50 companies by market cap, 2018-present (~60K rows)",
        "Target: Close (daily closing price)",
        "Merged with company fundamentals (sector, market cap, EBITDA)",
        "16 engineered features: lags, rolling averages, volume change",
    ])

    # 3. Feature correlation
    figure_slide(prs, "Feature Correlation", d, "figure_2.png",
                 "Strong linear relationships between Close and Open/High/Low")

    # 4. Models
    content_slide(prs, "Models Trained", [
        "Linear Regression (baseline)",
        "Lasso Regression (L1)",
        "Ridge Regression (L2)",
        "SVR (StandardScaler pipeline, RBF kernel)",
        "Decision Tree (max_depth=15)",
        "Random Forest (n_estimators=100, max_depth=20)",
        "Neural Network (Keras, 3 hidden layers, StandardScaler)",
    ])

    # 5. Hyperparameter tuning
    content_slide(prs, "Hyperparameter Tuning", [
        "Model tuned: Random Forest via GridSearchCV (3-fold CV)",
        "Parameters: n_estimators [100,200,300], max_depth [10,20,30], min_samples_split [2,5,10]",
        "Best: n_estimators=300, max_depth=30, min_samples_split=2",
        "Improvement: RMSE 1.0295 -> 1.0256 (0.37%)",
    ])

    # 6. Results table
    table_slide(prs, "Results", ["Model", "RMSE", "MAE", "R2"],
        [["Ridge", "0.0001", "0.0000", "1.0000"],
         ["Lasso", "0.0297", "0.0172", "1.0000"],
         ["Neural Network", "0.8693", "0.6073", "0.9999"],
         ["Random Forest (Tuned)", "1.0256", "0.5203", "0.9999"],
         ["Decision Tree", "1.4913", "0.7512", "0.9998"],
         ["Linear Regression", "2.4065", "1.4050", "0.9995"],
         ["SVR", "10.3737", "0.6864", "0.9914"]])

    # 7. Results chart
    figure_slide(prs, "Model Comparison", d, "figure_6.png",
                 "RMSE and R2 comparison across all models")

    # 8. Feature importance
    figure_slide(prs, "Feature Importances (Random Forest)", d, "figure_4.png",
                 "Low (59.1%) and High (40.6%) dominate")

    # 9. Key findings
    content_slide(prs, "Key Findings", [
        "All models achieved R2 > 0.99",
        "Lasso reduced features from 16 to 7 (only price features survived)",
        "Close is almost entirely determined by same-day OHLC features",
        "GridSearchCV gave modest improvement -- default RF was near-optimal",
    ])

    # 10. Conclusion
    content_slide(prs, "Conclusion", [
        "Recommended model: Ridge Regression",
        "Lowest RMSE (0.0001) and perfect R2 (1.0000)",
        "L2 regularization handles multicollinearity among correlated price features",
        "For forecasting future prices (no same-day data), tree-based models would likely outperform",
    ])

    prs.save(str(ROOT / d / 'part1_slides.pptx'))
    print('  part1_slides.pptx')


def build_part2():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    d = 'part2_classification'

    title_slide(prs, "Part 2: Classification", "Predicting S&P 500 Stock Price Direction")

    content_slide(prs, "Dataset", [
        "Same S&P 500 subset (top 50 companies, 2018-present)",
        "Target: binary -- 1 (UP) if next-day close > today, 0 (DOWN)",
        "Class balance: ~52% UP, ~48% DOWN",
        "13 technical indicator features: RSI, SMA ratio, momentum, volatility",
        "F1 Score as primary metric",
    ])

    figure_slide(prs, "Feature Correlation with Target", d, "figure_2.png",
                 "Weak correlations reflect randomness of daily price movements")

    content_slide(prs, "Models Trained", [
        "Logistic Regression (max_iter=1000)",
        "SVC (StandardScaler pipeline, RBF kernel)",
        "Decision Tree (max_depth=10)",
        "Random Forest (n_estimators=200, max_depth=15)",
        "Neural Network (Keras, 3 layers, sigmoid output, StandardScaler)",
    ])

    content_slide(prs, "Hyperparameter Tuning", [
        "Model tuned: Random Forest Classifier via GridSearchCV (3-fold, scoring=f1)",
        "Parameters: n_estimators [100,200,300], max_depth [10,15,20,None], min_samples_split [2,5,10]",
        "Best: n_estimators=300, max_depth=10, min_samples_split=10",
        "Improvement: F1 0.6638 -> 0.6778",
    ])

    table_slide(prs, "Results (sorted by F1 Score)",
        ["Model", "Accuracy", "F1", "Precision", "Recall"],
        [["Logistic Regression", "0.5277", "0.6829", "0.5288", "0.9636"],
         ["Random Forest (Tuned)", "0.5355", "0.6778", "0.5346", "0.9259"],
         ["SVC", "0.5284", "0.6730", "0.5307", "0.9196"],
         ["Decision Tree", "0.5293", "0.6689", "0.5320", "0.9009"],
         ["Neural Network", "0.5181", "0.5730", "0.5382", "0.6126"]])

    figure_slide(prs, "Accuracy vs F1 Score", d, "figure_5.png")

    figure_slide(prs, "Confusion Matrix (Tuned RF)", d, "figure_6.png",
                 "Slight bias toward predicting UP")

    content_slide(prs, "Key Findings & Conclusion", [
        "All models ~53% accuracy -- only slightly above random (50%)",
        "Logistic Regression: best F1 (0.6829), very high recall (0.9636)",
        "Daily stock direction is essentially a random walk with slight upward drift",
        "Complex models did not outperform logistic baseline",
        "Recommended: Logistic Regression",
    ])

    prs.save(str(ROOT / d / 'part2_slides.pptx'))
    print('  part2_slides.pptx')


def build_part3():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    d = 'part3_timeseries'

    title_slide(prs, "Part 3: Time Series", "Forecasting the S&P 500 Index")

    content_slide(prs, "Dataset", [
        "S&P 500 index daily values (Dec 2014 - Dec 2024)",
        "2,517 data points with genuine daily time index",
        "Target: S&P 500 index value (continuous, sequential)",
        "Chronological 80/20 split -- no shuffling",
    ])

    content_slide(prs, "Stationarity Analysis", [
        "ADF test on raw series: statistic=0.545, p=0.986 -- not stationary",
        "After first-order differencing: statistic=-15.941, p~0.0 -- stationary",
        "Volatility clusters visible (e.g., 2020 COVID crash)",
    ], part_dir=d, figure="figure_2.png")

    content_slide(prs, "Methods", [
        "ARIMA: manual grid search over (p,d,q), best = ARIMA(2,1,2), AIC=20,119",
        "Random Forest: 9 lag/rolling features (lag_1 to lag_7, rolling mean/std)",
        "Pipeline: StandardScaler -> Random Forest",
        "GridSearchCV with TimeSeriesSplit (3 splits)",
        "Best RF params: n_estimators=200, max_depth=20, min_samples_split=2",
    ])

    table_slide(prs, "Results",
        ["Method", "MAE", "RMSE", "Type"],
        [["Random Forest", "352.39", "544.17", "ML"],
         ["Random Forest (Tuned)", "352.39", "544.17", "ML"],
         ["ARIMA(2,1,2)", "1012.39", "1202.29", "Classical"]])

    figure_slide(prs, "Actual vs Forecasts", d, "figure_5.png",
                 "Random Forest tracks actual values closely; ARIMA diverges")

    figure_slide(prs, "MAE and RMSE Comparison", d, "figure_6.png")

    content_slide(prs, "Key Findings", [
        "Random Forest outperformed ARIMA -- 55% RMSE improvement",
        "ARIMA degrades over long horizons (projects flat trend)",
        "RF with lag features captures upward trajectory",
        "Tuned RF matched default -- parameters were already near-optimal",
    ])

    content_slide(prs, "Conclusion", [
        "Recommended: Random Forest with Lag Features",
        "RMSE improvement: 1202 -> 544 (55% better than ARIMA)",
        "ARIMA remains valuable for short-term (1-5 step) forecasts",
        "Combining both approaches could yield best practical results",
    ])

    prs.save(str(ROOT / d / 'part3_slides.pptx'))
    print('  part3_slides.pptx')


def build_part4():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    d = 'part4_clustering'

    title_slide(prs, "Part 4: Clustering", "Grouping S&P 500 Companies")

    content_slide(prs, "Dataset", [
        "11 numeric features: return, volatility, volume, market cap, EBITDA, etc.",
        "Fully unsupervised -- no target variable",
        "Combined daily stock aggregates with company fundamentals",
        "StandardScaler applied before all clustering algorithms",
        "PCA (2D) for visualization; clustering on full 11 features",
    ])

    figure_slide(prs, "PCA Projection", d, "figure_1.png",
                 "Mega-cap tech companies clearly separated from the main cluster")

    content_slide(prs, "Algorithms", [
        "K-Means: tested k=2 to 10, Elbow Method + Silhouette Scores",
        "DBSCAN: varied eps (0.5-5.0), min_samples (3,5,7,10)",
        "Agglomerative: 3 linkage methods (ward, complete, average) x k=2 to 5",
        "Selected best by Silhouette Score",
    ])

    figure_slide(prs, "K-Means: Elbow and Silhouette", d, "figure_2.png",
                 "Both indicate k=2 as optimal")

    table_slide(prs, "Clustering Comparison",
        ["Algorithm", "Clusters", "Silhouette", "Observation"],
        [["K-Means", "2", "0.7226", "Clean separation"],
         ["DBSCAN", "6", "0.3622", "131 noise points"],
         ["Agglomerative", "2", "0.8072", "Average linkage, best score"]])

    figure_slide(prs, "All Clustering Results", d, "figure_7.png",
                 "K-Means and Agglomerative produce similar 2-cluster solutions")

    content_slide(prs, "Key Findings", [
        "Natural 2-cluster split: mega-cap tech vs. rest of S&P 500",
        "Driven by Market Cap, EBITDA, and Index Weight",
        "DBSCAN identified 131 noise points (26% of companies)",
        "Dendrogram reveals hierarchical sub-group structures",
        "Average linkage with k=2 gave highest Silhouette (0.807)",
    ])

    content_slide(prs, "Conclusion", [
        "Recommended: Agglomerative Clustering (Average Linkage)",
        "Highest Silhouette Score (0.8072), cleanest 2-cluster solution",
        "Hierarchical structure provides dendrogram for sub-group analysis",
        "2-cluster solution useful for portfolio diversification",
        "Finer clustering (k=3-5) could separate growth vs. value vs. defensive",
    ])

    prs.save(str(ROOT / d / 'part4_slides.pptx'))
    print('  part4_slides.pptx')


if __name__ == '__main__':
    print('Generating slides...')
    build_part1()
    build_part2()
    build_part3()
    build_part4()
    print('Done.')
